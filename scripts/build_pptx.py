#!/usr/bin/env python3
"""
build_pptx.py — Build a PowerPoint presentation from Soapbox report data.

Usage:
    python build_pptx.py \\
        --template rsra \\
        --data /path/to/report_data.json \\
        --brand /path/to/brand.json \\
        --output /path/to/output.pptx \\
        [--templates-dir /path/to/templates]

    # --data and --brand also accept raw JSON strings or @/path/to/file.json:
    python build_pptx.py \\
        --template rsra \\
        --data '{"property": {...}, ...}' \\
        --brand '{"primary_color": "#1B2A3B"}' \\
        --output /path/to/report.pptx

Exits 0 on success (prints output path to stdout).
Exits 1 on error.
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from datetime import date
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
# Dependency bootstrap
# ---------------------------------------------------------------------------

def _ensure_pptx() -> None:
    """Install python-pptx if absent (handles Debian externally-managed envs)."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("python-pptx not found — installing...", file=sys.stderr)
        # Try normal install first; fall back to --break-system-packages on Debian/Ubuntu
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", "python-pptx"],
            capture_output=True,
        )
        if result.returncode != 0:
            subprocess.run(
                [sys.executable, "-m", "pip", "install", "python-pptx", "--break-system-packages"],
                check=True,
            )


_ensure_pptx()

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# Widescreen 16:9 slide dimensions (13.33" × 7.5")
SLIDE_WIDTH  = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Layout slots
HEADER_HEIGHT   = Inches(0.6)
FOOTER_HEIGHT   = Inches(0.25)
BODY_TOP        = HEADER_HEIGHT
BODY_HEIGHT     = SLIDE_HEIGHT - HEADER_HEIGHT - FOOTER_HEIGHT
MARGIN_H        = Inches(0.35)
CONTENT_WIDTH   = SLIDE_WIDTH - MARGIN_H * 2

# Signal level color map (hex strings)
SIGNAL_COLORS = {
    "Low Risk":                     "#2D6A4F",   # green
    "Moderate Risk — Opportunity":  "#B45309",   # amber
    "Moderate Risk — CapEx":        "#D97706",   # amber-orange
    "High Transition Risk":         "#B91C1C",   # red
}

ALTERNATING_ROW_FILL = RGBColor(0xF9, 0xFA, 0xFB)
TOTALS_ROW_FILL      = RGBColor(0xE5, 0xE7, 0xEB)


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a CSS hex color string to python-pptx RGBColor."""
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def _lighten(rgb: RGBColor, factor: float = 0.85) -> RGBColor:
    """Return a lightened version of an RGBColor (blend toward white)."""
    return RGBColor(
        int(rgb[0] + (255 - rgb[0]) * factor),
        int(rgb[1] + (255 - rgb[1]) * factor),
        int(rgb[2] + (255 - rgb[2]) * factor),
    )


# ---------------------------------------------------------------------------
# Brand resolution
# ---------------------------------------------------------------------------

DEFAULT_BRAND: dict[str, Any] = {
    "name": "Soapbox",
    "primary_color": "#0F1923",
    "highlight_color": "#52B788",
    "accent_color": "#F0F4F8",
    "text_color": "#0F1923",
    "text_muted": "#6B7280",
    "border_color": "#E5E7EB",
    "page_size": "Letter",
}


def resolve_brand(brand_arg: dict[str, Any] | None, templates_dir: Path | None, template: str) -> dict[str, Any]:
    """
    Merge brand data: defaults → org brand file (if resolvable) → brand_arg overrides.
    """
    merged = dict(DEFAULT_BRAND)

    # Attempt to load org brand from templates dir
    if templates_dir:
        # Try template-specific brand first, then soapbox default
        for candidate in [
            templates_dir / "_brand" / "soapbox" / "brand.json",
        ]:
            if candidate.exists():
                with open(candidate) as f:
                    merged.update(json.load(f))
                break

    if brand_arg:
        merged.update(brand_arg)

    return merged


# ---------------------------------------------------------------------------
# Slide builder helpers
# ---------------------------------------------------------------------------

def _blank_layout(prs: Presentation):
    """Return the blank slide layout (last in the default deck)."""
    return prs.slide_layouts[6]  # index 6 = Blank in most built-in themes


def _add_filled_rect(slide, left, top, width, height, fill_rgb: RGBColor):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text: str, *, font_size: int,
                  bold: bool = False, color: RGBColor | None = None,
                  align: PP_ALIGN = PP_ALIGN.LEFT, wrap: bool = True) -> Any:
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _add_footer(slide, brand: dict, slide_number: int, total_slides: int) -> None:
    """Add org name (left), date (center-ish), slide number (right) in muted 8pt."""
    muted = hex_to_rgb(brand.get("text_muted", "#6B7280"))
    y = SLIDE_HEIGHT - FOOTER_HEIGHT
    w_third = CONTENT_WIDTH // 3

    org = brand.get("name", "")
    today = date.today().isoformat()

    _add_text_box(slide, MARGIN_H, y, w_third, FOOTER_HEIGHT,
                  org, font_size=8, color=muted, align=PP_ALIGN.LEFT)
    _add_text_box(slide, MARGIN_H + w_third, y, w_third, FOOTER_HEIGHT,
                  today, font_size=8, color=muted, align=PP_ALIGN.CENTER)
    _add_text_box(slide, MARGIN_H + w_third * 2, y, w_third, FOOTER_HEIGHT,
                  f"{slide_number} / {total_slides}", font_size=8, color=muted, align=PP_ALIGN.RIGHT)


def _add_section_header(slide, title: str, brand: dict) -> None:
    """Full-width colored header bar with white title text."""
    primary = hex_to_rgb(brand["primary_color"])
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, HEADER_HEIGHT, primary)
    _add_text_box(slide, MARGIN_H, 0, SLIDE_WIDTH - MARGIN_H * 2, HEADER_HEIGHT,
                  title, font_size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                  align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Section renderers
# ---------------------------------------------------------------------------

def _render_table_section(slide, rows: list[list[str]], headers: list[str], brand: dict) -> None:
    """Render a table with branded headers and alternating row fills."""
    if not rows and not headers:
        return

    primary = hex_to_rgb(brand["primary_color"])
    n_cols = len(headers) if headers else len(rows[0])
    n_rows = len(rows) + (1 if headers else 0)

    # Determine column widths: equal distribution
    col_width = int(CONTENT_WIDTH / n_cols)
    row_height = Inches(0.28)
    table_top = BODY_TOP + Inches(0.15)
    max_rows_visible = int((BODY_HEIGHT - Inches(0.3)) / row_height)

    actual_rows = min(n_rows, max_rows_visible)
    table = slide.shapes.add_table(
        actual_rows, n_cols,
        MARGIN_H, table_top,
        CONTENT_WIDTH, row_height * actual_rows,
    ).table

    # Style helper
    def _cell_fill(cell, rgb: RGBColor):
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb

    def _cell_text(cell, text: str, font_size: int = 9, bold: bool = False,
                   color: RGBColor | None = None):
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(text) if text is not None else ""
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    row_offset = 0

    # Header row
    if headers:
        for ci, hdr in enumerate(headers[:n_cols]):
            cell = table.cell(0, ci)
            _cell_fill(cell, primary)
            _cell_text(cell, hdr, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        row_offset = 1

    # Data rows
    for ri, row_data in enumerate(rows[:actual_rows - row_offset]):
        tr = ri + row_offset
        is_totals = ri == len(rows) - 1 and len(rows) > 1
        fill = TOTALS_ROW_FILL if is_totals else (
            ALTERNATING_ROW_FILL if ri % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
        )
        for ci, val in enumerate(row_data[:n_cols]):
            cell = table.cell(tr, ci)
            _cell_fill(cell, fill)
            _cell_text(cell, val, bold=is_totals)


def _render_kv_section(slide, items: list[dict], brand: dict) -> None:
    """Render key-value pairs in two columns with a highlight left rule per row."""
    highlight = hex_to_rgb(brand.get("highlight_color", "#52B788"))
    muted = hex_to_rgb(brand.get("text_muted", "#6B7280"))
    primary_rgb = hex_to_rgb(brand["primary_color"])

    row_h = Inches(0.36)
    top = BODY_TOP + Inches(0.1)
    rule_w = Inches(0.04)
    label_w = Inches(2.5)
    value_w = CONTENT_WIDTH - label_w - rule_w - Inches(0.1)
    half = SLIDE_WIDTH // 2

    # Two-column layout: items split across left / right column
    left_items = items[: (len(items) + 1) // 2]
    right_items = items[(len(items) + 1) // 2 :]

    for col_items, col_x in [(left_items, MARGIN_H), (right_items, half + Inches(0.1))]:
        for ri, item in enumerate(col_items):
            y = top + ri * row_h
            label = item.get("label", "")
            value = item.get("value", "")

            # Highlight rule
            _add_filled_rect(slide, col_x, y + Inches(0.04), rule_w, row_h - Inches(0.08), highlight)

            # Label
            _add_text_box(slide, col_x + rule_w + Inches(0.06), y,
                          label_w, row_h, label,
                          font_size=10, bold=True, color=primary_rgb)

            # Value
            _add_text_box(slide, col_x + rule_w + Inches(0.06) + label_w, y,
                          value_w, row_h, str(value),
                          font_size=10, color=muted)


def _render_deal_signal(slide, data: dict, brand: dict) -> None:
    """Render a large color-coded signal box with level and narrative."""
    level = data.get("level", "")
    narrative = data.get("narrative", "")

    hex_color = SIGNAL_COLORS.get(level, brand.get("highlight_color", "#52B788"))
    signal_rgb = hex_to_rgb(hex_color)
    light_bg = _lighten(signal_rgb, 0.88)

    box_top = BODY_TOP + Inches(0.25)
    box_h = Inches(1.6)
    box_w = CONTENT_WIDTH

    # Background box
    _add_filled_rect(slide, MARGIN_H, box_top, box_w, box_h, light_bg)

    # Left accent bar
    _add_filled_rect(slide, MARGIN_H, box_top, Inches(0.12), box_h, signal_rgb)

    # Level text (large, bold, signal color)
    _add_text_box(slide, MARGIN_H + Inches(0.2), box_top + Inches(0.1),
                  box_w - Inches(0.3), Inches(0.5),
                  level, font_size=20, bold=True, color=signal_rgb)

    # Narrative text
    _add_text_box(slide, MARGIN_H + Inches(0.2), box_top + Inches(0.65),
                  box_w - Inches(0.3), Inches(0.85),
                  narrative, font_size=11, color=hex_to_rgb(brand.get("text_color", "#0F1923")))


def _render_narrative(slide, text: str, brand: dict) -> None:
    """Render a plain narrative text block below the header."""
    _add_text_box(slide, MARGIN_H, BODY_TOP + Inches(0.15),
                  CONTENT_WIDTH, BODY_HEIGHT - Inches(0.3),
                  text, font_size=11,
                  color=hex_to_rgb(brand.get("text_color", "#0F1923")))


# ---------------------------------------------------------------------------
# Section → slide dispatch
# ---------------------------------------------------------------------------

def _add_section_slide(prs: Presentation, title: str, section_data: Any,
                       section_key: str, brand: dict, slide_num: int, total: int) -> None:
    """Add a single section slide to the presentation."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _add_section_header(slide, title, brand)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Dispatch by section type / key
    if section_key == "deal_signal" and isinstance(section_data, dict):
        _render_deal_signal(slide, section_data, brand)

    elif section_key == "decarb_plan" and isinstance(section_data, list):
        headers = ["Measure", "Timing", "CapEx Total", "Incentive Program", "Financial Impact"]
        rows = []
        for m in section_data:
            rows.append([
                m.get("measure", ""),
                m.get("timing", ""),
                f"${m.get('capex_total', 0):,.0f}",
                m.get("incentive_program", "—"),
                m.get("financial_impact_value", "—"),
            ])
        # Totals row
        total_capex = sum(m.get("capex_total", 0) for m in section_data)
        rows.append(["TOTAL", "", f"${total_capex:,.0f}", "", ""])
        _render_table_section(slide, rows, headers, brand)

    elif section_key == "emissions_profile" and isinstance(section_data, dict):
        items = [
            {"label": "Fuel Profile",      "value": section_data.get("fuel_profile", "")},
            {"label": "Utility Structure", "value": section_data.get("utility_structure", "")},
            {"label": "Baseline Emissions","value": section_data.get("baseline_emissions", "")},
            {"label": "CRREM Pathway",     "value": section_data.get("crrem_pathway", "N/A")},
            {"label": "Regulation",        "value": section_data.get("regulation", "")},
            {"label": "EUI (kBtu/sqft)",   "value": section_data.get("eui_kbtu", "—")},
            {"label": "ENERGY STAR Score", "value": section_data.get("energy_star_score", "—")},
        ]
        items = [i for i in items if i["value"] not in ("", None, "—") or i["label"] in ("Regulation",)]
        _render_kv_section(slide, items, brand)

    elif section_key == "property" and isinstance(section_data, dict):
        items = [
            {"label": "Name",        "value": section_data.get("name", "")},
            {"label": "Address",     "value": section_data.get("address", "")},
            {"label": "Type",        "value": section_data.get("type", "").title()},
            {"label": "Year Built",  "value": section_data.get("year_built", "")},
            {"label": "Units",       "value": section_data.get("units", "—")},
            {"label": "GFA (sqft)",  "value": f'{section_data["gfa_sqft"]:,}' if section_data.get("gfa_sqft") else "—"},
            {"label": "ZIP",         "value": section_data.get("zip", "—")},
        ]
        _render_kv_section(slide, items, brand)

    elif section_key == "seller_questions" and isinstance(section_data, list):
        text = "\n".join(f"• {q}" for q in section_data)
        _render_narrative(slide, text, brand)

    elif isinstance(section_data, list):
        # Generic list → bullet narrative
        text = "\n".join(f"• {item}" for item in section_data if isinstance(item, str))
        _render_narrative(slide, text, brand)

    elif isinstance(section_data, dict):
        # Generic dict → key-value
        items = [{"label": k.replace("_", " ").title(), "value": str(v)}
                 for k, v in section_data.items() if v is not None]
        _render_kv_section(slide, items, brand)

    elif isinstance(section_data, str):
        _render_narrative(slide, section_data, brand)

    _add_footer(slide, brand, slide_num, total)


# ---------------------------------------------------------------------------
# Title slide
# ---------------------------------------------------------------------------

def _add_title_slide(prs: Presentation, data: dict, brand: dict, template: str) -> None:
    """Add a full-bleed branded title slide."""
    slide = prs.slides.add_slide(_blank_layout(prs))

    primary = hex_to_rgb(brand["primary_color"])
    highlight = hex_to_rgb(brand.get("highlight_color", "#52B788"))
    white = RGBColor(0xFF, 0xFF, 0xFF)

    # Full-bleed background
    _add_filled_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, primary)

    # Accent stripe
    stripe_h = Inches(0.06)
    _add_filled_rect(slide, 0, SLIDE_HEIGHT - stripe_h - Inches(0.5), SLIDE_WIDTH, stripe_h, highlight)

    # Property name
    prop = data.get("property", {})
    prop_name = prop.get("name", "Report")
    _add_text_box(slide, MARGIN_H, Inches(2.0), SLIDE_WIDTH - MARGIN_H * 2, Inches(1.2),
                  prop_name, font_size=40, bold=True, color=white)

    # Report type label
    template_labels = {
        "rsra":              "Rapid Sustainability Risk Assessment",
        "retrofit-plan":     "Decarbonization Capital Plan",
        "crrem-assessment":  "CRREM Pathway Assessment",
        "bps-compliance":    "Building Performance Standards Analysis",
        "decarb-roadmap":    "Decarbonization Roadmap",
        "portfolio-summary": "Portfolio ESG Summary",
        "gresb-submission":  "GRESB Reporting Package",
    }
    label = template_labels.get(template, template.replace("-", " ").title())
    _add_text_box(slide, MARGIN_H, Inches(3.4), SLIDE_WIDTH - MARGIN_H * 2, Inches(0.5),
                  label, font_size=16, color=highlight)

    # Date and org
    report_date = data.get("report_date", date.today().isoformat())
    org_line = brand.get("name", "")
    prepared_for = data.get("prepared_for", "")
    if prepared_for:
        org_line = f"{org_line}  ·  {prepared_for}"
    subtitle = f"{report_date}  ·  {org_line}" if org_line else report_date
    _add_text_box(slide, MARGIN_H, Inches(4.1), SLIDE_WIDTH - MARGIN_H * 2, Inches(0.4),
                  subtitle, font_size=12, color=white)


# ---------------------------------------------------------------------------
# Section ordering
# ---------------------------------------------------------------------------

SECTION_TITLES = {
    "property":          "Property Overview",
    "emissions_profile": "Emissions Profile",
    "decarb_plan":       "Decarbonization Capital Plan",
    "deal_signal":       "Deal Signal",
    "seller_questions":  "Seller Questions",
    "data_quality":      "Data Quality",
}

DEFAULT_SECTION_ORDER = [
    "deal_signal",
    "property",
    "emissions_profile",
    "decarb_plan",
    "seller_questions",
    "data_quality",
]

SKIP_KEYS = {
    "report_date", "prepared_by", "prepared_for", "disposition_mode",
    "decarb_plan_total",
}


def _section_order(data: dict) -> list[str]:
    """Return section keys in the preferred display order."""
    if "section_order" in data:
        return [k for k in data["section_order"] if k in data and k not in SKIP_KEYS]

    ordered = [k for k in DEFAULT_SECTION_ORDER if k in data and k not in SKIP_KEYS]
    extras = [k for k in data if k not in ordered and k not in SKIP_KEYS
              and k not in DEFAULT_SECTION_ORDER]
    return ordered + extras


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build_pptx(
    template: str,
    data: dict,
    brand_overrides: dict | None,
    output_path: Path,
    templates_dir: Path | None,
) -> Path:
    """Build and save the PPTX. Returns the output path."""
    brand = resolve_brand(brand_overrides, templates_dir, template)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    sections = _section_order(data)
    # +1 for title slide; pre-compute total for footer
    total_slides = 1 + len(sections)

    _add_title_slide(prs, data, brand, template)

    for idx, key in enumerate(sections, start=1):
        title = SECTION_TITLES.get(key, key.replace("_", " ").title())
        _add_section_slide(prs, title, data[key], key, brand, idx + 1, total_slides)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _load_json_arg(raw: str) -> dict:
    """Load JSON from a string, a @filepath reference, or a plain file path."""
    # Strip @ prefix (legacy filepath syntax)
    if raw.startswith("@"):
        raw = raw[1:]
    # If it looks like a file path (exists on disk or ends in .json), treat as file
    candidate = Path(raw)
    try:
        if candidate.exists() and candidate.is_file():
            return json.loads(candidate.read_text())
    except OSError:
        pass  # Path too long — fall through to JSON literal parse
    # Otherwise parse as a JSON literal
    return json.loads(raw)


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Build a PowerPoint from Soapbox report data.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    p.add_argument("--template", required=True, help="Report template name (e.g. rsra)")
    p.add_argument(
        "--data",
        required=True,
        help="Path to report_data.json, a JSON file path, or a raw JSON string",
    )
    p.add_argument(
        "--brand",
        default=None,
        help="Path to brand.json, a JSON file path, or a raw JSON string (optional)",
    )
    p.add_argument("--output", required=True, help="Destination .pptx path")
    p.add_argument(
        "--templates-dir",
        default=None,
        help="Path to the templates/ directory (for brand resolution)",
    )
    return p


def main() -> None:
    args = _build_parser().parse_args()

    try:
        data = _load_json_arg(args.data)
    except json.JSONDecodeError as e:
        print(f"Invalid JSON in --data: {e}", file=sys.stderr)
        sys.exit(1)

    brand_overrides: dict | None = None
    if args.brand:
        try:
            brand_overrides = _load_json_arg(args.brand)
        except json.JSONDecodeError as e:
            print(f"Invalid JSON in --brand: {e}", file=sys.stderr)
            sys.exit(1)

    output_path = Path(args.output).resolve()
    templates_dir = Path(args.templates_dir).resolve() if args.templates_dir else None

    try:
        result = build_pptx(
            template=args.template,
            data=data,
            brand_overrides=brand_overrides,
            output_path=output_path,
            templates_dir=templates_dir,
        )
        print(str(result))
    except Exception as exc:
        print(f"build_pptx error: {exc}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
